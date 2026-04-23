"""Build the single final-project presentation slide.

Produces one polished 16:9 slide with:
- Title
- Key takeaway (one-liner)
- 3 concise bullets
- The headline figure (regime-dependence heatmap)

Output: Quant_Final_Slide.pptx
"""

from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from scipy import stats

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

sns.set_style("whitegrid")

TICKERS = ["SPY", "QQQ", "XLK", "XLE", "XLF", "XLU"]
HORIZONS = [1, 5, 20]

# --- Recompute the same tables used in the report ---------------------------
prices = pd.read_csv("data/etf_prices.csv", index_col=0, parse_dates=True)
prices = prices[TICKERS].sort_index()
prices.index = pd.to_datetime(prices.index)
prices = prices.reindex(pd.bdate_range(prices.index.min(), prices.index.max())).ffill(limit=1).dropna()
returns = np.log(prices).diff().dropna()

features = {}
for t in TICKERS:
    r = returns[t]
    df = pd.DataFrame({"ret": r})
    df["vol_20"] = r.rolling(20).std() * np.sqrt(252)
    for h in HORIZONS:
        df[f"lag_{h}"] = r.rolling(h).sum()
        df[f"fwd_{h}"] = r.shift(-1).rolling(h).sum().shift(-(h - 1))
    features[t] = df

spy_vol = features["SPY"]["vol_20"].dropna()
q25 = spy_vol.expanding(252).quantile(0.25)
q75 = spy_vol.expanding(252).quantile(0.75)
regime = pd.Series("Mid Vol", index=spy_vol.index)
regime[spy_vol <= q25] = "Low Vol"
regime[spy_vol >= q75] = "High Vol"

def pearson_by_regime(reg_label):
    out = {}
    for t in TICKERS:
        merged = features[t].join(regime.rename("reg"))
        out[t] = {}
        for h in HORIZONS:
            sub = merged[merged["reg"] == reg_label][[f"lag_{h}", f"fwd_{h}"]].dropna()
            out[t][h] = stats.pearsonr(sub[f"lag_{h}"], sub[f"fwd_{h}"]).statistic
    return pd.DataFrame(out).T.reindex(TICKERS)[HORIZONS]

low_r = pearson_by_regime("Low Vol")
high_r = pearson_by_regime("High Vol")

# --- Build the headline figure (saved as PNG for embedding) ----------------
FIG_PATH = Path("data/slide_headline.png")

fig, axes = plt.subplots(1, 2, figsize=(11, 4.4))
sns.heatmap(low_r, annot=True, fmt=".3f", cmap="RdBu_r", center=0,
            vmin=-0.25, vmax=0.25, ax=axes[0], cbar=False)
axes[0].set_title("Low-Volatility Regime", fontsize=12)
axes[0].set_xlabel("Horizon (trading days)")
axes[0].set_ylabel("ETF")

sns.heatmap(high_r, annot=True, fmt=".3f", cmap="RdBu_r", center=0,
            vmin=-0.25, vmax=0.25, ax=axes[1], cbar_kws={"shrink": 0.8})
axes[1].set_title("High-Volatility Regime", fontsize=12)
axes[1].set_xlabel("Horizon (trading days)")
axes[1].set_ylabel("")

fig.suptitle(
    "Pearson correlation of past h-day return vs. next h-day return (2010-2026)",
    fontsize=12, fontweight="bold", y=1.02,
)
plt.tight_layout()
plt.savefig(FIG_PATH, dpi=180, bbox_inches="tight")
plt.close(fig)

# --- Build the pptx slide --------------------------------------------------
SLIDE_PATH = Path("Quant_Final_Slide.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

# Background accent bar along the top.
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.4))
bar.fill.solid()
bar.fill.fore_color.rgb = RGBColor(0x0B, 0x3D, 0x91)
bar.line.fill.background()

def add_text(left, top, width, height, text, *, size=14, bold=False,
             color=RGBColor(0x1A, 0x1A, 0x1A), italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

# Title
add_text(Inches(0.5), Inches(0.55), Inches(12.3), Inches(0.7),
         "Mean Reversion in US Sector ETFs Depends on Volatility Regime",
         size=28, bold=True, color=RGBColor(0x0B, 0x3D, 0x91))

# Subtitle / takeaway
add_text(Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.5),
         "Short-horizon return reversal is mild in calm markets and ~3-4x stronger in high-volatility markets.",
         size=16, italic=True, color=RGBColor(0x33, 0x33, 0x33))

# Figure
slide.shapes.add_picture(str(FIG_PATH), Inches(0.5), Inches(1.9),
                         width=Inches(8.5), height=Inches(3.9))

# Bullets box on the right
bullets_tb = slide.shapes.add_textbox(Inches(9.2), Inches(1.9), Inches(3.9), Inches(4.5))
tf = bullets_tb.text_frame
tf.word_wrap = True

def add_bullet(tf, text, *, bold_lead=None, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(6)
    if bold_lead:
        r1 = p.add_run(); r1.text = bold_lead + " "
        r1.font.bold = True; r1.font.size = Pt(13)
        r1.font.color.rgb = RGBColor(0x0B, 0x3D, 0x91)
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(13)
    r2.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

add_bullet(tf, "Every ETF shows negative lagged/forward correlation at h = 1 day.",
           bold_lead="Pervasive 1d reversal.", first=True)
add_bullet(tf, "At h = 1d, Low-Vol r is near zero; High-Vol r ranges from -0.10 to -0.17.",
           bold_lead="Regime amplifies reversal.")
add_bullet(tf, "XLU (utilities) has the strongest 20-day reversal (r = -0.22 full-sample, -0.29 high-vol).",
           bold_lead="Sector dispersion matters.")
add_bullet(tf, "XLE (energy) is the only ETF with a positive 5-day correlation — commodity momentum, not equity reversal.",
           bold_lead="One outlier.")

# Footer takeaway box
footer_tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9))
tf = footer_tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r1 = p.add_run(); r1.text = "Takeaway.  "
r1.font.bold = True; r1.font.size = Pt(15); r1.font.color.rgb = RGBColor(0x0B, 0x3D, 0x91)
r2 = p.add_run()
r2.text = ("Return-predictability patterns in US equity ETFs are not uniform — they depend jointly on "
           "horizon, sector, and volatility regime. Conditional analysis makes a mildly-noisy full-sample "
           "average look like a cleanly-structured signal.")
r2.font.size = Pt(13); r2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Footer meta
add_text(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
         "Data: Yahoo Finance (SPY, QQQ, XLK, XLE, XLF, XLU), 2010-01-05 to 2026-04-17, 4,248 trading days. "
         "Regime: SPY 20d realized-vol quartiles (expanding).",
         size=10, italic=True, color=RGBColor(0x66, 0x66, 0x66))

prs.save(SLIDE_PATH)
print(f"Wrote {SLIDE_PATH}")
